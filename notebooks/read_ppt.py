from pptx import Presentation

prs = Presentation(r"notebooks/Template PPT Final Project ET4243.pptx")
for i, slide in enumerate(prs.slides):
    layout = slide.slide_layout.name
    print(f"=== Slide {i+1} ({layout}) ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
            if texts:
                print(f"  [{shape.name}]")
                for t in texts[:4]:
                    print(f"    - {t[:100]}")
    print()
